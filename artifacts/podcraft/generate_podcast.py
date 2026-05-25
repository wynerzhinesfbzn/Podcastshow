import os
import json
import base64
import subprocess
import tempfile
import threading
import requests
from pathlib import Path
from io import BytesIO
from openai import OpenAI

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
    base_url=os.environ.get("OPENAI_BASE_URL"),
)

STYLE_PROMPTS = {
    "anime": "anime art style, vibrant colors, expressive characters, Japanese animation aesthetic, dynamic composition",
    "cinematic": "cinematic film still, dramatic lighting, photorealistic, movie scene, professional cinematography",
    "cartoon": "cartoon illustration style, bold outlines, bright cheerful colors, fun and playful, comic book aesthetic",
    "futurism": "futuristic digital art, neon colors, cyberpunk aesthetic, sci-fi technology, glowing elements, dark background",
    "infographic": "clean infographic design, flat vector art, minimal style, geometric shapes, professional data visualization",
}

BASE_MUSIC = "https://incompetech.com/music/royalty-free/mp3-royaltyfree"
MUSIC_POOL = [
    {
        "id": "thinking",
        "name": "💭 Музыка для размышлений",
        "mood": "Спокойная, задумчивая — идеал для подкаста",
        "url": f"{BASE_MUSIC}/Thinking%20Music.mp3",
    },
    {
        "id": "lobby",
        "name": "☕ Лаундж-атмосфера",
        "mood": "Мягкий лаундж, уютный фон",
        "url": f"{BASE_MUSIC}/Lobby%20Time.mp3",
    },
    {
        "id": "piano",
        "name": "🎹 Спокойное фортепиано",
        "mood": "Тихое соло-фортепиано, очень нейтральное",
        "url": f"{BASE_MUSIC}/Relaxing%20Piano%20Music.mp3",
    },
    {
        "id": "bossa",
        "name": "🎷 Босса-нова / Джаз",
        "mood": "Лёгкий джаз, босса-нова, живой ритм",
        "url": f"{BASE_MUSIC}/Bossa%20Antigua.mp3",
    },
    {
        "id": "chill",
        "name": "😌 Чилл-аут",
        "mood": "Расслабленный, очень медленный темп",
        "url": f"{BASE_MUSIC}/Chill.mp3",
    },
    {
        "id": "backed_vibes",
        "name": "🎧 Чилл-хоп",
        "mood": "Мягкий хип-хоп бит, современный подкаст-стиль",
        "url": f"{BASE_MUSIC}/Backed%20Vibes%20Clean.mp3",
    },
    {
        "id": "slow_burn",
        "name": "🌊 Медленный эмбиент",
        "mood": "Атмосферный, глубокий, почти без ритма",
        "url": f"{BASE_MUSIC}/Slow%20Burn.mp3",
    },
    {
        "id": "meditation",
        "name": "🧘 Медитативный",
        "mood": "Минималистичный, тихий, без отвлечений",
        "url": f"{BASE_MUSIC}/Meditation%20Impromptu%2001.mp3",
    },
    {
        "id": "beauty_flow",
        "name": "✨ Нежный поток",
        "mood": "Мягкий и красивый, лёгкая мелодия",
        "url": f"{BASE_MUSIC}/Beauty%20Flow.mp3",
    },
    {
        "id": "peaceful",
        "name": "🌅 Тихая атмосфера",
        "mood": "Очень спокойный, почти тишина с мелодией",
        "url": f"{BASE_MUSIC}/Peaceful%20Desolation.mp3",
    },
    {
        "id": "jazz_groove",
        "name": "🕵️ Джазовый грув",
        "mood": "Игривый джаз, чуть бодрее — для живых тем",
        "url": f"{BASE_MUSIC}/Sneaky%20Snitch.mp3",
    },
    {
        "id": "decisions",
        "name": "🎬 Кинематографический фон",
        "mood": "Лёгкий кино-стиль, нейтральный и серьёзный",
        "url": f"{BASE_MUSIC}/Decisions.mp3",
    },
]

MUSIC_TRACKS = MUSIC_POOL  # backward compat — full pool exposed

ALL_VOICES = [
    {"id": "alloy",   "label": "Alloy",   "desc": "Нейтральный, универсальный", "gender": "neutral"},
    {"id": "ash",     "label": "Ash",     "desc": "Мягкий, спокойный",          "gender": "neutral"},
    {"id": "ballad",  "label": "Ballad",  "desc": "Певучий, тёплый",             "gender": "neutral"},
    {"id": "coral",   "label": "Coral",   "desc": "Женский, дружелюбный",        "gender": "female"},
    {"id": "echo",    "label": "Echo",    "desc": "Мужской, чёткий",             "gender": "male"},
    {"id": "fable",   "label": "Fable",   "desc": "Выразительный, театральный",  "gender": "neutral"},
    {"id": "nova",    "label": "Nova",    "desc": "Женский, живой",              "gender": "female"},
    {"id": "onyx",    "label": "Onyx",    "desc": "Глубокий, басовитый",         "gender": "male"},
    {"id": "sage",    "label": "Sage",    "desc": "Умеренный, профессиональный", "gender": "neutral"},
    {"id": "shimmer", "label": "Shimmer", "desc": "Тёплый, лёгкий",             "gender": "female"},
]

VALID_VOICE_IDS = {v["id"] for v in ALL_VOICES}

VOICE_PERSONALITIES = [
    {
        "id": "deep_thinker",
        "label": "🧠 Глубокий аналитик",
        "voice": "onyx",
        "short": "Вдумчивый, философский",
        "style": (
            "Говорит медленно и вдумчиво, делает паузы для размышления, строит глубокие рассуждения. "
            "Любит риторические вопросы. Типичные фразы: «Знаете, я об этом долго думал...», "
            "«Если посмотреть глубже...», «Это интересный парадокс», «А вот что меня поражает...»."
        ),
    },
    {
        "id": "radio_host",
        "label": "🎙 Радиоведущий",
        "voice": "alloy",
        "short": "Чёткий, профессиональный",
        "style": (
            "Говорит чётко и уверенно, как опытный радиоведущий. Структурирует мысли, плавные переходы, "
            "хорошо держит темп беседы. Типичные фразы: «Итак, давайте разберём...», «Переходим к следующей теме», "
            "«Это ключевой момент, друзья», «И вот почему это важно»."
        ),
    },
    {
        "id": "emotional",
        "label": "🎭 Эмоциональный",
        "voice": "nova",
        "short": "Живой, реактивный, эмоции",
        "style": (
            "Реагирует эмоционально и непосредственно, живые реакции, часто восклицает, выражает искреннее "
            "удивление и восторг. Типичные фразы: «Вот это да!», «Подождите, это же означает...», "
            "«Я честно не ожидал!», «Это просто невероятно!»."
        ),
    },
    {
        "id": "warm_expert",
        "label": "😊 Добродушный эксперт",
        "voice": "shimmer",
        "short": "Тёплый, дружелюбный, с юмором",
        "style": (
            "Тёплый и дружелюбный тон, объясняет сложное простыми словами, иногда шутит и смеётся. "
            "Типичные фразы: «Представьте себе такую картину...», «Говоря проще...», "
            "«И тут начинается самое интересное!», «Ну это же просто класс, согласитесь?»."
        ),
    },
    {
        "id": "energetic",
        "label": "⚡ Энергичный",
        "voice": "ash",
        "short": "Быстрый, задорный, заряженный",
        "style": (
            "Говорит энергично, задаёт темп беседе, быстро переходит от мысли к мысли, заряжает аудиторию. "
            "Типичные фразы: «Отлично, поехали!», «Это работает вот так — слушайте внимательно!», "
            "«Давайте дальше, это только разогрев!»."
        ),
    },
    {
        "id": "dry_business",
        "label": "🧊 Деловой",
        "voice": "echo",
        "short": "Сухой, фактический, по делу",
        "style": (
            "Говорит чётко и лаконично, только факты без воды. Минимум эмоций, максимум конкретики. "
            "Типичные фразы: «Факт: ...», «По данным исследований...», «Вывод очевиден», «Следующий пункт»."
        ),
    },
    {
        "id": "charismatic",
        "label": "🌟 Харизматичный",
        "voice": "coral",
        "short": "Обаятельный, убедительный",
        "style": (
            "Харизматичный и убедительный, рассказывает истории, вовлекает слушателя, умеет держать паузу. "
            "Типичные фразы: «Вот вам история...», «И знаете, что самое удивительное?», "
            "«Это меняет всё», «Представьте — вы просыпаетесь и...»."
        ),
    },
    {
        "id": "theatrical",
        "label": "🎭 Театральный",
        "voice": "fable",
        "short": "Выразительный, образный",
        "style": (
            "Говорит образно и выразительно, как актёр. Использует яркие метафоры, делает акценты. "
            "Типичные фразы: «Представьте такую картину...», «Это как если бы...», "
            "«Занавес открывается...», «Вот где начинается настоящая драма»."
        ),
    },
]

NARRATOR_PRESETS = [
    {
        "id": "pro_anchor",
        "label": "📢 Профессиональный диктор",
        "voice": "sage",
        "short": "Чёткий, нейтральный, TV-качество",
        "style": (
            "Профессиональный теле/радиодиктор. Чёткая дикция, нейтральный тон, максимальная разборчивость. "
            "Представляет гостей объективно и без лишних слов. "
            "Никогда не участвует в дискуссии — только представляет, переходит, закрывает эфир."
        ),
    },
    {
        "id": "solemn",
        "label": "🏆 Торжественный",
        "voice": "onyx",
        "short": "Глубокий, авторитетный, весомый",
        "style": (
            "Глубокий, басовитый голос. Торжественные объявления, каждое слово весомо. "
            "Представляет гостей с достоинством. Идеален для серьёзных тем и документальных форматов."
        ),
    },
    {
        "id": "warm_anchor",
        "label": "🌟 Тёплый анонсер",
        "voice": "coral",
        "short": "Дружелюбный, приятный, вовлекающий",
        "style": (
            "Тёплый и располагающий голос диктора. Представляет гостей с симпатией, "
            "создаёт атмосферу дружеской передачи. Плавные переходы между темами."
        ),
    },
    {
        "id": "dynamic_anchor",
        "label": "⚡ Динамичный",
        "voice": "alloy",
        "short": "Энергичный, радийный, современный",
        "style": (
            "Энергичный современный радиоведущий-диктор. Динамичные переходы, цепкие анонсы. "
            "Идеален для молодёжных и развлекательных форматов."
        ),
    },
]

LANGUAGES = [
    {"id": "ru",   "label": "🇷🇺 Русский",     "name_en": "Russian"},
    {"id": "en",   "label": "🇺🇸 English",      "name_en": "English"},
    {"id": "es",   "label": "🇪🇸 Español",      "name_en": "Spanish"},
    {"id": "fr",   "label": "🇫🇷 Français",     "name_en": "French"},
    {"id": "de",   "label": "🇩🇪 Deutsch",      "name_en": "German"},
    {"id": "it",   "label": "🇮🇹 Italiano",     "name_en": "Italian"},
    {"id": "pt",   "label": "🇧🇷 Português",    "name_en": "Portuguese"},
    {"id": "zh",   "label": "🇨🇳 中文",          "name_en": "Chinese"},
    {"id": "ja",   "label": "🇯🇵 日本語",        "name_en": "Japanese"},
    {"id": "ko",   "label": "🇰🇷 한국어",        "name_en": "Korean"},
    {"id": "ar",   "label": "🇸🇦 العربية",      "name_en": "Arabic"},
    {"id": "hi",   "label": "🇮🇳 हिन्दी",       "name_en": "Hindi"},
    {"id": "tr",   "label": "🇹🇷 Türkçe",       "name_en": "Turkish"},
    {"id": "pl",   "label": "🇵🇱 Polski",       "name_en": "Polish"},
    {"id": "nl",   "label": "🇳🇱 Nederlands",   "name_en": "Dutch"},
]


def update_progress(job_data: dict, step: str, percent: int, message: str):
    job_data["step"] = step
    job_data["percent"] = percent
    job_data["message"] = message


def _get_lang_name(language: str) -> str:
    lang = next((l for l in LANGUAGES if l["id"] == language), None)
    return lang["name_en"] if lang else "Russian"


def generate_dialogue(
    text: str,
    duration_minutes: int,
    language: str = "ru",
    host_name: str = "Host",
    guest_name: str = "Guest",
    host_style: str = "",
    guest_style: str = "",
    narrator_name: str = "Диктор",
    narrator_style: str = "",
    host_description: str = "",
    guest_description: str = "",
    update_cb=None,
) -> list[dict]:
    words_per_minute = 130
    target_words = duration_minutes * words_per_minute
    lang_name = _get_lang_name(language)

    host_style_block = f"Speaking style: {host_style}" if host_style else "Speaking style: confident expert, sets the tone and pace"
    guest_style_block = f"Speaking style: {guest_style}" if guest_style else "Speaking style: curious and engaged, asks clarifying questions and reacts naturally"
    narrator_style_block = narrator_style if narrator_style else "Professional, neutral TV/radio announcer voice. Clear, concise, authoritative."

    host_desc_line = f" — {host_description}" if host_description else ""
    guest_desc_line = f" — {guest_description}" if guest_description else ""

    prompt = f"""You are an elite professional radio podcast scriptwriter. Create an outstanding, radio-quality podcast script with THREE roles: a NARRATOR/ANNOUNCER, and two hosts.

CRITICAL: Write the ENTIRE script in {lang_name}. Every single word must be in {lang_name}. Zero exceptions.

═══ ROLES ═══

NARRATOR — "{narrator_name}":
- Style: {narrator_style_block}
- ONLY appears at: show opening, segment transitions, show closing
- NEVER discusses, debates, or shares opinions on the topic
- Introduces hosts and guest by name with a brief description
- Speaks in 3rd person about the hosts ("Слово {host_name}...", "Сегодня в студии — {guest_name}...")
- Keep each narrator line SHORT: 1–2 sentences max
- Uses: segment number 0 for intro, highest segment number for outro, intermediate segments for transitions

HOST A — "{host_name}"{host_desc_line}:
- {host_style_block}
- Uses their natural speech patterns consistently

HOST B — "{guest_name}"{guest_desc_line}:
- {guest_style_block}
- Uses their natural speech patterns consistently

═══ SCRIPT STRUCTURE ═══

segment 0 — NARRATOR opens the show (1–2 sentences): introduces show topic, then "{host_name}" and "{guest_name}" by name with short descriptions
segment 0 — HOST A and HOST B start discussing (natural conversation, engaging hook)
segments 1–N — HOST A and HOST B discuss the topic (bulk of content, ~{target_words} words total)
  — Between each segment: NARRATOR says 1 sentence transition ("Переходим к следующей теме...", etc.)
final segment — NARRATOR closes the show (1–2 sentences, memorable)

═══ DIALOGUE RULES ═══
1. Target TOTAL length: ~{target_words} words (~{duration_minutes} minutes of spoken audio)
2. Sound like a REAL radio conversation — natural flow, NOT a lecture or Q&A
3. Hosts address each other by name occasionally but naturally
4. Include natural speech elements: pauses ("..."), reactions ("ага", "да-да"), gentle interruptions
5. NO robotic transitions in host lines. Each line must feel spoken, not written
6. Distribute host speaking time roughly equally
7. Divide into 4–10 thematic segments (segment: 0, 1, 2, ...)
8. NARRATOR lines use segment 0 for intro, highest number for outro, in-between for transitions

═══ RESPONSE FORMAT — strict JSON only ═══
{{
  "dialogue": [
    {{"speaker": "narrator", "text": "opening line in {lang_name}", "segment": 0}},
    {{"speaker": "host", "text": "spoken line in {lang_name}", "segment": 0}},
    {{"speaker": "guest", "text": "spoken line in {lang_name}", "segment": 0}},
    {{"speaker": "narrator", "text": "transition in {lang_name}", "segment": 1}},
    {{"speaker": "host", "text": "spoken line in {lang_name}", "segment": 1}},
    {{"speaker": "guest", "text": "spoken line in {lang_name}", "segment": 1}},
    {{"speaker": "narrator", "text": "closing line in {lang_name}", "segment": 9}}
  ]
}}

Source material to transform into a podcast:
{text[:8000]}"""

    if update_cb:
        update_cb("dialogue", 10, "GPT генерирует диалог подкаста...")

    response = client.chat.completions.create(
        model="gpt-5.4",
        messages=[{"role": "user", "content": prompt}],
        max_completion_tokens=6000,
        response_format={"type": "json_object"},
    )

    raw = response.choices[0].message.content
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        import re
        match = re.search(r'\[.*\]', raw, re.DOTALL)
        if match:
            data = json.loads(match.group(0))
            return data
        raise ValueError("GPT вернул некорректный JSON")

    if isinstance(data, dict):
        if "dialogue" in data and isinstance(data["dialogue"], list):
            return data["dialogue"]
        for key in data:
            if isinstance(data[key], list) and len(data[key]) > 0:
                return data[key]
    if isinstance(data, list):
        return data

    raise ValueError(f"Неожиданный формат ответа GPT: {str(data)[:200]}")


def _tts_standard(text: str, safe_voice: str, lang_name: str) -> bytes:
    """Стандартная TTS без клонирования голоса."""
    response = client.chat.completions.create(
        model="gpt-audio",
        modalities=["text", "audio"],
        audio={"voice": safe_voice, "format": "mp3"},
        messages=[
            {
                "role": "system",
                "content": (
                    f"You are a world-class professional broadcast voice talent for premium radio and podcast production. "
                    f"You speak flawless, native-level {lang_name} — perfectly clear, zero accent, indistinguishable from a top human broadcaster. "
                    "Your delivery is warm, natural, and deeply engaging — exactly like the hosts on premium audio shows such as Google NotebookLM or Wondery productions. "
                    "Speak the user's text word for word, exactly as written — do NOT add, skip, or alter any words. "
                    "Apply these professional voice techniques: "
                    "— Natural breathing rhythm with micro-pauses at punctuation marks; "
                    "— Warm emotional tone that invites the listener in; "
                    "— Varied pace: slow down on key points, speed up on transitions; "
                    "— Perfect pronunciation of every word with confident articulation; "
                    "— No robotic flatness — genuine human warmth and energy in every sentence; "
                    "— Where '...' appears, insert a natural thoughtful pause; "
                    f"— Accent, intonation, and stress patterns must be 100% native {lang_name}."
                ),
            },
            {"role": "user", "content": text[:4096]},
        ],
        max_completion_tokens=8192,
    )
    return base64.b64decode(response.choices[0].message.audio.data)


def _tts_cloned(text: str, safe_voice: str, lang_name: str, voice_sample_path: str) -> bytes:
    """TTS с клонированием голоса из загруженного аудио-образца."""
    from pydub import AudioSegment

    seg = AudioSegment.from_file(voice_sample_path)
    seg = seg[:30000]
    seg = seg.set_frame_rate(16000).set_channels(1)
    wav_buf = BytesIO()
    seg.export(wav_buf, format="wav")
    sample_b64 = base64.b64encode(wav_buf.getvalue()).decode("utf-8")

    response = client.chat.completions.create(
        model="gpt-audio",
        modalities=["text", "audio"],
        audio={"voice": safe_voice, "format": "mp3"},
        messages=[
            {
                "role": "system",
                "content": (
                    f"You are an advanced AI voice cloning and synthesis system. "
                    f"You will receive an audio sample of a real person's voice, then text to speak. "
                    f"Your task: reproduce the EXACT voice characteristics of the person in the audio sample — "
                    f"their unique timbre, pitch, speaking pace, accent, intonation patterns, and emotional tone. "
                    f"The output must be 100% in {lang_name}. "
                    f"Speak the given text word for word, exactly as written, but in the cloned voice. "
                    f"Do NOT add or remove any words. "
                    f"Apply natural breathing, micro-pauses at punctuation, and warm human delivery."
                ),
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "input_audio",
                        "input_audio": {
                            "data": sample_b64,
                            "format": "wav",
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            f"Speak in the exact voice from the audio sample above. "
                            f"Say this text word for word: {text[:3800]}"
                        ),
                    },
                ],
            },
        ],
        max_completion_tokens=8192,
    )
    return base64.b64decode(response.choices[0].message.audio.data)


def generate_tts(
    text: str,
    voice: str,
    output_path: str,
    language: str = "ru",
    voice_sample_path: str | None = None,
):
    """
    Генерирует MP3 через gpt-audio.
    Если voice_sample_path задан — клонирует голос из образца.
    """
    safe_voice = voice if voice in VALID_VOICE_IDS else "alloy"
    lang_name = _get_lang_name(language)

    if voice_sample_path and os.path.exists(voice_sample_path):
        try:
            audio_bytes = _tts_cloned(text, safe_voice, lang_name, voice_sample_path)
        except Exception:
            audio_bytes = _tts_standard(text, safe_voice, lang_name)
    else:
        audio_bytes = _tts_standard(text, safe_voice, lang_name)

    with open(output_path, "wb") as f:
        f.write(audio_bytes)


def merge_audio_with_ducking(speech_files: list[str], music_url: str, output_path: str):
    """Склеивает речевые MP3-файлы и накладывает фоновую музыку с ducking-эффектом."""
    from pydub import AudioSegment

    if not speech_files:
        silence = AudioSegment.silent(duration=3000)
        silence.export(output_path, format="mp3")
        return

    combined = AudioSegment.empty()
    for f in speech_files:
        try:
            seg = AudioSegment.from_file(f)
            combined += seg + AudioSegment.silent(duration=250)
        except Exception:
            continue

    if len(combined) == 0:
        AudioSegment.silent(duration=3000).export(output_path, format="mp3")
        return

    try:
        music_resp = requests.get(music_url, timeout=30)
        music_resp.raise_for_status()
        music_raw = BytesIO(music_resp.content)
        music = AudioSegment.from_file(music_raw, format="mp3")

        total_ms = len(combined)
        while len(music) < total_ms:
            music = music + music
        music = music[:total_ms]

        ducked_music = AudioSegment.empty()
        pos = 0
        for f in speech_files:
            try:
                seg = AudioSegment.from_file(f)
            except Exception:
                continue
            dur = len(seg)
            ducked_music += (music - 28)[pos: pos + dur]
            pos += dur
            pause_dur = 250
            ducked_music += (music - 16)[pos: pos + pause_dur]
            pos += pause_dur

        if len(ducked_music) < len(combined):
            ducked_music += AudioSegment.silent(duration=len(combined) - len(ducked_music))

        final = combined.overlay(ducked_music)
    except Exception:
        final = combined

    final.export(output_path, format="mp3")


def generate_image_for_segment(segment_text: str, style: str, output_path: str):
    """Генерирует PNG через gpt-image-1 (DALL-E 3 compatible)."""
    style_hint = STYLE_PROMPTS.get(style, STYLE_PROMPTS["cinematic"])
    prompt = (
        f"Illustration for a podcast episode. Topic: {segment_text[:300]}. "
        f"Visual style: {style_hint}. "
        "Widescreen composition, no text, no labels, no watermarks."
    )

    response = client.images.generate(
        model="gpt-image-1",
        prompt=prompt,
        size="1024x1024",
        n=1,
    )

    item = response.data[0]
    b64 = getattr(item, "b64_json", None)
    url = getattr(item, "url", None)

    if b64:
        with open(output_path, "wb") as f:
            f.write(base64.b64decode(b64))
    elif url:
        img_data = requests.get(url, timeout=30).content
        with open(output_path, "wb") as f:
            f.write(img_data)
    else:
        _make_placeholder_image(output_path, segment_text)


def _make_placeholder_image(output_path: str, text: str = ""):
    """Создаёт заглушку-изображение если AI недоступен."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (1024, 1024), color=(20, 20, 40))
    draw = ImageDraw.Draw(img)
    draw.rectangle([40, 40, 984, 984], outline=(100, 100, 200), width=3)
    draw.text((512, 480), "🎙 PodCraft", fill=(150, 150, 255), anchor="mm")
    if text:
        short = text[:80] + ("..." if len(text) > 80 else "")
        draw.text((512, 550), short, fill=(180, 180, 220), anchor="mm")
    img.save(output_path)


def create_mp4(image_paths: list[str], durations: list[float], audio_path: str, output_path: str):
    """Собирает анимированный MP4 с Ken Burns эффектами (zoom/pan) через FFmpeg."""
    if not image_paths:
        raise ValueError("Нет изображений для сборки видео")

    FPS = 25
    PATTERNS = [
        # zoom in center
        ("min(zoom+0.0015,1.5)", "iw/2-(iw/zoom/2)", "ih/2-(ih/zoom/2)"),
        # pan right + slight zoom
        ("1.2",                  "min(iw*0.2,0.4*on)", "ih/2-(ih/zoom/2)"),
        # zoom out from center
        ("max(1.0,1.5-0.0012*on)", "iw/2-(iw/zoom/2)", "ih/2-(ih/zoom/2)"),
        # pan left + slight zoom
        ("1.2",                  "max(0,iw*0.2-0.4*on)", "ih/2-(ih/zoom/2)"),
        # pan up + zoom
        ("1.2",                  "iw/2-(iw/zoom/2)", "max(0,ih*0.15-0.3*on)"),
    ]

    tmp_dir = tempfile.mkdtemp()
    clip_files: list[str] = []

    try:
        for idx, (img, dur) in enumerate(zip(image_paths, durations)):
            dur = max(dur, 3.0)
            nframes = int(dur * FPS) + 5
            z_expr, x_expr, y_expr = PATTERNS[idx % len(PATTERNS)]
            zp = (
                f"zoompan=z='{z_expr}':x='{x_expr}':y='{y_expr}'"
                f":d={nframes}:s=1024x1024:fps={FPS}"
            )
            clip_path = os.path.join(tmp_dir, f"clip_{idx:04d}.mp4")
            r = subprocess.run(
                ["ffmpeg", "-y", "-loop", "1", "-i", img,
                 "-vf", zp, "-t", f"{dur:.3f}",
                 "-c:v", "libx264", "-pix_fmt", "yuv420p", "-an", clip_path],
                capture_output=True, text=True,
            )
            if r.returncode != 0:
                # Fallback: static clip
                subprocess.run(
                    ["ffmpeg", "-y", "-loop", "1", "-i", img,
                     "-vf", "scale=1024:1024:force_original_aspect_ratio=decrease,"
                            "pad=1024:1024:(ow-iw)/2:(oh-ih)/2:color=black",
                     "-t", f"{dur:.3f}", "-r", str(FPS),
                     "-c:v", "libx264", "-pix_fmt", "yuv420p", "-an", clip_path],
                    capture_output=True,
                )
            clip_files.append(clip_path)

        concat_file = os.path.join(tmp_dir, "concat.txt")
        with open(concat_file, "w") as f:
            for cp in clip_files:
                f.write(f"file '{cp}'\n")

        video_only = output_path.replace(".mp4", "_noaudio.mp4")
        r = subprocess.run(
            ["ffmpeg", "-y", "-f", "concat", "-safe", "0",
             "-i", concat_file, "-c", "copy", video_only],
            capture_output=True, text=True,
        )
        if r.returncode != 0:
            raise RuntimeError(f"FFmpeg concat: {r.stderr[-500:]}")

        r = subprocess.run(
            ["ffmpeg", "-y", "-i", video_only, "-i", audio_path,
             "-c:v", "copy", "-c:a", "aac", "-b:a", "192k", "-shortest", output_path],
            capture_output=True, text=True,
        )
        if r.returncode != 0:
            raise RuntimeError(f"FFmpeg mux: {r.stderr[-500:]}")

    finally:
        for f in clip_files:
            try: os.unlink(f)
            except Exception: pass
        for name in ["concat.txt", "podcast_noaudio.mp4"]:
            try: os.unlink(os.path.join(tmp_dir, name))
            except Exception: pass
        video_only_path = output_path.replace(".mp4", "_noaudio.mp4")
        if os.path.exists(video_only_path):
            try: os.unlink(video_only_path)
            except Exception: pass
        try: os.rmdir(tmp_dir)
        except Exception: pass


def create_pptx(image_paths: list[str], segments_text: list[str], output_path: str):
    """Создаёт PPTX-презентацию с изображениями и тезисами."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0F, 0x11, 0x17)
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "🎙 PodCraft — Подкаст"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x81, 0x8C, 0xF8)
    p.alignment = PP_ALIGN.CENTER

    for i, (img_path, text) in enumerate(zip(image_paths, segments_text)):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1D, 0x2E)

        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(7.5), Inches(7.5))

        txBox = slide.shapes.add_textbox(Inches(7.7), Inches(0.4), Inches(5.3), Inches(6.8))
        tf = txBox.text_frame
        tf.word_wrap = True

        p_num = tf.add_paragraph()
        p_num.text = f"СЕГМЕНТ {i + 1} / {len(segments_text)}"
        p_num.font.size = Pt(10)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(0x7C, 0x83, 0xED)

        p_text = tf.add_paragraph()
        p_text.text = text[:350]
        p_text.font.size = Pt(15)
        p_text.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        p_text.space_before = Pt(10)

    prs.save(output_path)


def extract_segments(dialogue: list[dict]) -> list[tuple[int, str]]:
    """
    Возвращает список (segment_id, combined_text) отсортированный по segment_id.
    """
    segments: dict[int, list[str]] = {}
    for line in dialogue:
        seg = int(line.get("segment", 0))
        segments.setdefault(seg, []).append(line.get("text", ""))

    result = []
    for k in sorted(segments.keys()):
        combined = " ".join(segments[k])
        result.append((k, combined[:400]))
    return result


def generate_podcast(
    text: str,
    style: str,
    host_voice: str,
    guest_voice: str,
    duration_minutes: int,
    music_id: str,
    output_dir: str,
    job_data: dict,
    language: str = "ru",
    podcast_mode: str = "video",   # "audio" | "video"
    host_name: str = "Ведущий",
    guest_name: str = "Гость",
    narrator_voice: str = "sage",
    narrator_name: str = "Диктор",
    host_sample: str | None = None,
    guest_sample: str | None = None,
    narrator_sample: str | None = None,
):
    os.makedirs(output_dir, exist_ok=True)
    speech_dir = os.path.join(output_dir, "speech")
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(speech_dir, exist_ok=True)
    os.makedirs(images_dir, exist_ok=True)

    def cb(step, pct, msg):
        update_progress(job_data, step, pct, msg)

    # ── Шаг A: Генерация диалога ─────────────────────────────────────────
    cb("dialogue", 5, "GPT генерирует диалог подкаста...")
    dialogue = generate_dialogue(
        text, duration_minutes,
        language=language,
        host_name=host_name,
        guest_name=guest_name,
        host_style=job_data.get("host_style", ""),
        guest_style=job_data.get("guest_style", ""),
        narrator_name=narrator_name,
        narrator_style=job_data.get("narrator_style", ""),
        host_description=job_data.get("host_description", ""),
        guest_description=job_data.get("guest_description", ""),
        update_cb=cb,
    )
    if not dialogue:
        raise ValueError("GPT вернул пустой диалог")
    job_data["dialogue"] = dialogue
    cb("dialogue", 20, f"Диалог готов: {len(dialogue)} реплик")

    # ── Шаг B: TTS для каждой реплики ────────────────────────────────────
    cb("tts", 21, "Озвучка реплик через OpenAI TTS...")
    from pydub import AudioSegment

    speech_files: list[str] = []
    segment_durations: dict[int, float] = {}
    speech_timeline: list[dict] = []   # [{speaker, name, start_s, end_s, text}]
    tts_errors: list[str] = []
    current_ms = 0.0
    GAP_MS = 250  # пауза между репликами

    for i, line in enumerate(dialogue):
        speaker = line.get("speaker", "host")
        if speaker == "narrator":
            name = narrator_name
            voice = narrator_voice
        elif speaker == "guest":
            name = guest_name
            voice = guest_voice
        else:
            name = host_name
            voice = host_voice
        text_line = (line.get("text") or "").strip()
        seg_num = int(line.get("segment", 0))

        if not text_line:
            continue

        out_file = os.path.join(speech_dir, f"line_{i:04d}.mp3")
        sample_path = narrator_sample if speaker == "narrator" else (guest_sample if speaker == "guest" else host_sample)
        try:
            generate_tts(text_line, voice, out_file, language=language, voice_sample_path=sample_path)
            seg = AudioSegment.from_file(out_file)
            dur_ms = len(seg)
            speech_files.append(out_file)
            speech_timeline.append({
                "speaker": speaker,
                "name": name,
                "start_s": round(current_ms / 1000, 3),
                "end_s": round((current_ms + dur_ms) / 1000, 3),
                "text": text_line[:200],
            })
            segment_durations[seg_num] = segment_durations.get(seg_num, 0.0) + dur_ms / 1000.0
            current_ms += dur_ms + GAP_MS
        except Exception as e:
            tts_errors.append(f"Реплика {i}: {e}")

        pct = 21 + int((i + 1) / max(len(dialogue), 1) * 29)
        cb("tts", min(pct, 49), f"Озвучка реплики {i + 1}/{len(dialogue)}...")

    if tts_errors:
        job_data["tts_errors"] = tts_errors

    if not speech_files:
        raise ValueError("Не удалось озвучить ни одной реплики.")

    job_data["speech_timeline"] = speech_timeline

    # ── Шаг C: Склейка аудио с музыкой + ducking ─────────────────────────
    cb("audio", 50, "Склейка аудио с фоновой музыкой...")
    music_track = next((t for t in MUSIC_TRACKS if t["id"] == music_id), MUSIC_TRACKS[0])
    mp3_path = os.path.join(output_dir, "podcast.mp3")
    merge_audio_with_ducking(speech_files, music_track["url"], mp3_path)
    cb("audio", 57, f"MP3 готов! ({os.path.getsize(mp3_path) // 1024} КБ)")

    mp4_path = os.path.join(output_dir, "podcast.mp4")
    pptx_path = os.path.join(output_dir, "podcast.pptx")
    mp4_ok = False
    pptx_ok = False
    image_paths: list[str] = []
    image_durations: list[float] = []

    if podcast_mode == "video":
        # ── Шаги D+E: Генерация изображений + анимированное видео ─────────
        cb("images", 58, "Генерация изображений через DALL-E 3...")
        segments_data = extract_segments(dialogue)
        image_errors: list[str] = []

        total_audio_sec = sum(segment_durations.values()) or (duration_minutes * 60)
        default_img_dur = total_audio_sec / max(len(segments_data), 1)

        for idx, (seg_id, seg_text) in enumerate(segments_data):
            img_path = os.path.join(images_dir, f"segment_{idx:03d}.png")
            try:
                generate_image_for_segment(seg_text, style, img_path)
            except Exception as e:
                image_errors.append(f"Сегмент {idx}: {e}")
                _make_placeholder_image(img_path, seg_text)

            image_paths.append(img_path)
            dur = segment_durations.get(seg_id, default_img_dur)
            image_durations.append(max(dur, 3.0))

            pct = 58 + int((idx + 1) / max(len(segments_data), 1) * 22)
            cb("images", min(pct, 79), f"Изображение {idx + 1}/{len(segments_data)}...")

        if image_errors:
            job_data["image_errors"] = image_errors

        cb("video", 80, "Сборка анимированного MP4 через FFmpeg...")
        try:
            create_mp4(image_paths, image_durations, mp3_path, mp4_path)
            mp4_ok = os.path.exists(mp4_path) and os.path.getsize(mp4_path) > 1024
            if mp4_ok:
                cb("video", 90, f"MP4 готов! ({os.path.getsize(mp4_path) // 1024} КБ)")
            else:
                cb("video", 90, "MP4: файл не создан")
        except Exception as e:
            job_data.setdefault("errors", []).append(f"MP4: {e}")
            cb("video", 90, f"MP4 ошибка: {str(e)[:80]}")

        # ── Шаг F: PPTX ──────────────────────────────────────────────────
        cb("pptx", 91, "Создание презентации PPTX...")
        segments_text_only = [t for _, t in segments_data]
        try:
            create_pptx(image_paths, segments_text_only, pptx_path)
            pptx_ok = os.path.exists(pptx_path) and os.path.getsize(pptx_path) > 1024
            if pptx_ok:
                cb("pptx", 98, f"PPTX готов! ({os.path.getsize(pptx_path) // 1024} КБ)")
        except Exception as e:
            job_data.setdefault("errors", []).append(f"PPTX: {e}")
            cb("pptx", 98, f"PPTX ошибка: {str(e)[:80]}")
    else:
        segments_text_only = []

    # ── Сохраняем результаты ──────────────────────────────────────────────
    job_data["mp3_path"] = mp3_path if os.path.exists(mp3_path) else None
    job_data["mp4_path"] = mp4_path if mp4_ok else None
    job_data["pptx_path"] = pptx_path if pptx_ok else None
    job_data["music_tracks"] = MUSIC_TRACKS
    job_data["image_paths"] = image_paths
    job_data["image_durations"] = image_durations
    job_data["segment_durations"] = segment_durations
    job_data["segments_text"] = segments_text_only
    job_data["podcast_mode"] = podcast_mode
    job_data["host_name"] = host_name
    job_data["guest_name"] = guest_name
    job_data["status"] = "done"
    cb("done", 100, "Подкаст успешно создан!")
